import os
import shutil
import json
import unittest
import pandas as pd
from unittest.mock import patch
from pptx import Presentation
from pptx.util import Inches

import clean_data
import generate_pitch
import deck_refresh
import ib_check_deck
import model_update

class TestOtherScripts(unittest.TestCase):
    def setUp(self):
        self.test_dir = "./out/test_other_data"
        if os.path.exists(self.test_dir):
            try:
                shutil.rmtree(self.test_dir)
            except PermissionError:
                import time
                time.sleep(0.5)
                shutil.rmtree(self.test_dir, ignore_errors=True)
        os.makedirs(self.test_dir, exist_ok=True)

        self.ticker = "7203"
        self.ticker_dir = os.path.join(self.test_dir, "7203.T_Toyota_Motor")
        self.market_data_dir = os.path.join(self.ticker_dir, "market_data")
        os.makedirs(self.market_data_dir, exist_ok=True)

        # 財務データ（get_latest_financial_data用）を準備
        summary = {
            "long_name": "Toyota Motor",
            "ticker": "7203.T",
            "currency": "JPY",
            "current_price": 2700.0,
            "shares_outstanding": 11841052480,
            "market_cap": 32000000000000,
            "ebitda": 5000000000000
        }
        with open(os.path.join(self.market_data_dir, "summary.json"), "w", encoding="utf-8") as f:
            json.dump(summary, f)

        inc_df = pd.DataFrame({
            "2026-03-31": [40000.0, 5000.0, 4500.0, 500.0, 100.0]
        }, index=["Total Revenue", "EBITDA", "EBIT", "Reconciled Depreciation", "Tax Provision"])
        inc_df.to_csv(os.path.join(self.market_data_dir, "annual_income_stmt.csv"))

        bs_df = pd.DataFrame({
            "2026-03-31": [80000.0, 10000.0, 9000.0]
        }, index=["Total Assets", "Total Debt", "Cash Cash Equivalents And Short Term Investments"])
        bs_df.to_csv(os.path.join(self.market_data_dir, "annual_balance_sheet.csv"))

        cf_df = pd.DataFrame({
            "2026-03-31": [3000.0, 2000.0, 100.0]
        }, index=["Operating Cash Flow", "Capital Expenditure", "Change In Working Capital"])
        cf_df.to_csv(os.path.join(self.market_data_dir, "annual_cashflow.csv"))

        # prices.csv（クレンジングテスト用）
        prices_df = pd.DataFrame({
            "Date": [" 2026-06-19 00:00:00+09:00", " 2026-06-20 00:00:00+09:00", " 2026-06-20 00:00:00+09:00"], # 空白・重複・タイムゾーン
            "Close": [" 2700.0", "2776.5 ", "2776.5 "], # スペース・欠損値で重複行にする
            "Open": [2680.0, 2750.0, 2750.0],
            "High": [2710.0, 2790.0, 2790.0],
            "Low": [2670.0, 2740.0, 2740.0],
            "Volume": [1000, 2000, 2000]
        })
        prices_df.to_csv(os.path.join(self.market_data_dir, "prices.csv"), index=False)

    def tearDown(self):
        if os.path.exists(self.test_dir):
            try:
                shutil.rmtree(self.test_dir)
            except PermissionError:
                import time
                time.sleep(0.5)
                shutil.rmtree(self.test_dir, ignore_errors=True)

    @patch('sys.argv', ['clean_data.py', '7203', '--outdir', './out/test_other_data'])
    def test_clean_data(self):
        clean_data.main()
        cleaned_path = os.path.join(self.market_data_dir, "prices.csv")
        self.assertTrue(os.path.exists(cleaned_path))
        
        df = pd.read_csv(cleaned_path)
        self.assertEqual(len(df), 2) # 重複排除で2行になるはず
        self.assertEqual(df.loc[0, "Date"], "2026-06-19") # タイムゾーン・空白正規化
        self.assertEqual(df.loc[1, "Close"], 2776.5) # トリム・数値化

    @patch('sys.argv', ['generate_pitch.py', '7203', '--outdir', './out/test_other_data'])
    def test_generate_pitch(self):
        generate_pitch.main()
        
        analysis_dir = os.path.join(self.ticker_dir, "analysis")
        pptx_files = os.listdir(analysis_dir)
        self.assertTrue(any(f.endswith(".pptx") for f in pptx_files))

    @patch('sys.argv', ['deck_refresh.py', '7203', '--replacements', 'Toyota:Honda', '--outdir', './out/test_other_data'])
    def test_deck_refresh(self):
        analysis_dir = os.path.join(self.ticker_dir, "analysis")
        os.makedirs(analysis_dir, exist_ok=True)
        pptx_path = os.path.join(analysis_dir, "test_deck.pptx")

        # ダミーのPPTXを作成
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        tx_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(10), Inches(1))
        tf = tx_box.text_frame
        tf.paragraphs[0].text = "This is a Toyota presentation."
        prs.save(pptx_path)

        # deck_refresh.py を、--pptx-path を引数に追加して実行
        with patch('sys.argv', ['deck_refresh.py', '7203', '--pptx-path', pptx_path, '--replacements', 'Toyota:Honda', '--outdir', './out/test_other_data']):
            deck_refresh.main()

        # 置換後の検証
        updated_prs = Presentation(pptx_path)
        updated_text = updated_prs.slides[0].shapes[0].text_frame.text
        self.assertIn("Honda", updated_text)

    @patch('sys.argv', ['ib_check_deck.py', '7203', '--outdir', './out/test_other_data'])
    def test_ib_check_deck(self):
        analysis_dir = os.path.join(self.ticker_dir, "analysis")
        os.makedirs(analysis_dir, exist_ok=True)
        pptx_path = os.path.join(analysis_dir, "pitch_7203.T.pptx")

        # テスト用の表を含むPPTXを作成
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # 表を追加 (4行 x 2列)
        rows, cols = 4, 2
        left, top, width, height = Inches(1), Inches(1), Inches(6), Inches(2)
        table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # セル値を書き込み (Revenue)
        table.cell(0, 0).text = "Revenue"
        table.cell(0, 1).text = "40000"
        prs.save(pptx_path)

        # ib_check_deck.main() 呼び出し
        with patch('sys.argv', ['ib_check_deck.py', '7203', '--pptx-path', pptx_path, '--outdir', './out/test_other_data']):
            ib_check_deck.main()

    @patch('sys.argv', ['model_update.py', '7203', '--revenue', '50000', '--ebit', '5000', '--ebitda', '10000', '--outdir', './out/test_other_data'])
    def test_model_update(self):
        # model_update.main() 実行
        model_update.main()

        # CSVのバックアップができているか確認 (タイムスタンプ付き)
        import glob
        bak_files = glob.glob(os.path.join(self.market_data_dir, "annual_income_stmt.csv.*.bak"))
        self.assertTrue(len(bak_files) > 0)

        # 監査履歴ファイル (update_history.json) が作成されているか確認
        history_file = os.path.join(self.market_data_dir, "update_history.json")
        self.assertTrue(os.path.exists(history_file))
        with open(history_file, "r", encoding="utf-8") as f:
            history = json.load(f)
            self.assertTrue(len(history) > 0)
            self.assertIn("revenue", history[0]["changes"])

        # 値が更新されているか確認
        df = pd.read_csv(os.path.join(self.market_data_dir, "annual_income_stmt.csv"), index_col=0)
        self.assertEqual(df.loc["Total Revenue"].iloc[0], 50000.0)

    def test_model_update_partial_none(self):
        # 一部の値が None（指定なし）の場合のテスト
        with patch('sys.argv', ['model_update.py', '7203', '--revenue', '60000', '--outdir', './out/test_other_data']):
            model_update.main()
        df = pd.read_csv(os.path.join(self.market_data_dir, "annual_income_stmt.csv"), index_col=0)
        self.assertEqual(df.loc["Total Revenue"].iloc[0], 60000.0)
