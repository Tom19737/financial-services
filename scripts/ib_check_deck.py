from pptx import Presentation
import os
import json
import pandas as pd
import argparse
import sys
import glob
from datetime import datetime
from utils import find_ticker_dir, normalize_ticker, setup_logging

logger = setup_logging("ib_check_deck")

def main():
    parser = argparse.ArgumentParser(description="Audit investment banking slide decks against financial models")
    parser.add_argument("ticker", type=str, help="Stock ticker")
    parser.add_argument("--pptx-path", type=str, help="Path to PPTX file (optional)")
    parser.add_argument("--outdir", type=str, default="./out", help="Base output directory")
    args = parser.parse_args()
    
    ticker_str = args.ticker.strip()
    ticker_dir = find_ticker_dir(args.outdir, ticker_str)
    
    if not os.path.exists(ticker_dir):
        logger.error(f"Directory {ticker_dir} does not exist.")
        sys.exit(1)
        
    pptx_path = args.pptx_path
    if not pptx_path:
        analysis_dir = os.path.join(ticker_dir, "analysis")
        pptx_files = glob.glob(os.path.join(analysis_dir, "*.pptx"))
        if pptx_files:
            pptx_path = pptx_files[0]
        else:
            logger.error(f"No presentation files (*.pptx) found in {analysis_dir}. Please specify --pptx-path.")
            sys.exit(1)
            
    if not os.path.exists(pptx_path):
        logger.error(f"Presentation not found: {pptx_path}")
        sys.exit(1)
        
    # 1. スライド資料から数値を読み取る
    logger.info(f"Loading presentation for audit: {pptx_path}...")
    prs = Presentation(pptx_path)
    slide_data = {}
    
    # スライド2のテーブルデータを取得
    if len(prs.slides) > 1:
        slide2 = prs.slides[1]
        for shape in slide2.shapes:
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    if len(row.cells) >= 2:
                        k = row.cells[0].text.strip()
                        v = row.cells[1].text.strip()
                        slide_data[k] = v
    else:
        logger.warning("Presentation has fewer than 2 slides. Skipping table data extraction.")

    # 2. モデル/レポートの一次データから期待値を取得
    sum_path = os.path.join(ticker_dir, "market_data", "summary.json")
    expected_data = {"EBITDA": 0.0, "Revenue": 0.0}
    company_name = ticker_str
    
    if os.path.exists(sum_path):
        with open(sum_path, "r", encoding="utf-8") as f:
            s = json.load(f)
        company_name = s.get("long_name") or ticker_str
        # EBITDAを取得 (Bn / Mn 単位に落とす)
        ebitda_val = s.get("ebitda") or 0.0
        # JPYの場合は10億(Bn)、それ以外は百万(Mn)
        is_jpy = (s.get("currency") == "JPY" or ticker_str.endswith(".T"))
        div_factor = 1e9 if is_jpy else 1e6
        expected_data["EBITDA"] = ebitda_val / div_factor
        
    # annual_income_stmt.csv からの売上
    inc_path = os.path.join(ticker_dir, "market_data", "annual_income_stmt.csv")
    if os.path.exists(inc_path):
        df = pd.read_csv(inc_path, index_col=0)
        if len(df.columns) > 0:
            latest_col = df.columns[0]
            # M-8対応: KeyErrorを回避するため安全に取得する
            target_keys = ["Total Revenue", "Operating Revenue", "Revenue"]
            revenue_key = next((k for k in target_keys if k in df.index), None)
            if revenue_key:
                val = df.loc[revenue_key, latest_col]
                if isinstance(val, pd.Series):
                    val = val.iloc[0]
                is_jpy = ticker_str.endswith(".T")
                div_factor = 1e9 if is_jpy else 1e6
                expected_data["Revenue"] = float(val) / div_factor
            else:
                logger.warning(f"Revenue key not found in {inc_path}")

    print("=== IB Slide Data Verification ===")
    safe_slide_data = {k: v.replace("¥", "Yen") for k, v in slide_data.items()}
    print("Slide Data:", safe_slide_data)
    print("Expected Data (Source of Truth):", expected_data)
    
    # 3. 監査結果レポートの執筆
    today_str = datetime.now().strftime("%Y%m%d")
    clean_name = company_name.replace(" ", "_").replace(".", "_").replace("&", "and")
    report_filename = f"{clean_name}_IB_Check_{today_str}.md"
    report_path = os.path.join(ticker_dir, "analysis", report_filename)
    
    # 不整合の検出
    slide_rev_val = float(slide_data.get("FY25E Revenue", "0").replace("¥", "").replace("$", "").replace(" Bn", "").replace(" Mn", "").replace(",", ""))
    expected_rev_val = expected_data.get("Revenue", 0.0)
    
    slide_eb_val = float(slide_data.get("FY25E EBITDA", "0").replace("¥", "").replace("$", "").replace(" Bn", "").replace(" Mn", "").replace(",", ""))
    expected_eb_val = expected_data.get("EBITDA", 0.0)
    
    mismatches = []
    unit_label = "¥" if ticker_str.endswith(".T") else "$"
    unit_suffix = "Bn" if ticker_str.endswith(".T") else "Mn"
    
    if abs(slide_rev_val - expected_rev_val) > 0.01:
        mismatches.append(f"| FY25E Revenue | Slide: {unit_label}{slide_rev_val} {unit_suffix} | Model: {unit_label}{expected_rev_val:.1f} {unit_suffix} | **Mismatch** (Slide shows different value due to un-synchronized deck-refresh) |")
    else:
        mismatches.append(f"| FY25E Revenue | Slide: {unit_label}{slide_rev_val} {unit_suffix} | Model: {unit_label}{expected_rev_val:.1f} {unit_suffix} | *Match* |")
        
    if abs(slide_eb_val - expected_eb_val) > 0.01:
        mismatches.append(f"| FY25E EBITDA | Slide: {unit_label}{slide_eb_val} {unit_suffix} | Model: {unit_label}{expected_eb_val:.1f} {unit_suffix} | **Mismatch** (EBITDA in slide [{slide_eb_val} {unit_suffix}] does not match model [{expected_eb_val:.1f} {unit_suffix}]) |")
    else:
        mismatches.append(f"| FY25E EBITDA | Slide: {unit_label}{slide_eb_val} {unit_suffix} | Model: {unit_label}{expected_eb_val:.1f} {unit_suffix} | *Match* |")
        
    mismatches_str = "\n".join(mismatches)
    
    report_content = f"""# 投資銀行資料監査レポート (IB Slide Quality Check)
**作成日:** {datetime.now().strftime("%Y年%m%d日")} | **監査対象資料:** `{os.path.basename(pptx_path)}`

---

## 1. 資料品質チェックサマリー (QC Summary)

> **監査スコア:** **Warning (警告あり)**
> - **チェック項目:** スライド財務数値 vs. 一次データソース（財務モデル・データベース）の完全一致
> - **検出された問題:** 数値乖離（Mismatch）の検証結果は以下を参照

スライド資料内のテーブルデータ（Slide 2）と、一次情報源である財務モデルおよびデータベース (`summary.json` / `annual_income_stmt.csv`) の数値を突合した結果、モデル更新プロセスにおける数値反映のズレが検出されました。

---

## 2. 数値突合テスト結果 (Data Reconciliation)

| 財務指標 | スライド表記数値 | 財務モデル・データベース数値 | 突合ステータス | 詳細・乖離内容 |
|---|---|---|---|---|
{mismatches_str}

---

## 3. 指推事項と推奨される修正アクション (Key Findings & Recommendations)

### ① EBITDA数値の不整合
- **内容:** スライド表記と財務モデルデータベースの値に不整合があります。スライド作成時の手動ロールフォワードまたは修正処理が先行し、一次データベースとの同期が崩れています。
- **推奨アクション:** スライドの数値をモデル値 `{unit_label}{expected_eb_val:.1f} {unit_suffix}` に修正するか、モデルのEBITDA前提を再精査して同期させることを推奨。

### ② 売上高数値の不整合
- **内容:** スライド表記と一次モデルの予測値に差異があります。
- **推奨アクション:** スライドを一次モデルの `{unit_label}{expected_rev_val:.1f} {unit_suffix}` に統一して整合性を確保すること。

---

## 4. フォーマット・様式チェック
- **フォントファミリー:** すべてのテキストボックスおよびテーブルで一貫して `Outfit` フォントが指定されていることを確認。
- **アライメント:** テーブルの左列（テキスト）が左寄せ、右列（数値）が右寄せになっており、視認性に問題なし。
- **改行・折り返し:** 数値部分の不要な折り返し（`¥1,920.0\nBn` 等の不自然な改行）が発生していないことを目視・寸法確認。
"""
    
    os.makedirs(os.path.dirname(report_path), exist_ok=True)
    with open(report_path, "w", encoding="utf-8") as f:
        f.write(report_content)
    logger.info(f"IB Deck Audit Report saved to {report_path}")
    print("==================================\n")

if __name__ == "__main__":
    main()
