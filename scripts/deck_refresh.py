from pptx import Presentation
import os
import argparse
import sys
import glob
from utils import find_ticker_dir, normalize_ticker, setup_logging

logger = setup_logging("deck_refresh")

def refresh_text(shape, old_text, new_text):
    if shape.has_text_frame:
        tf = shape.text_frame
        for paragraph in tf.paragraphs:
            for run in paragraph.runs:
                if old_text in run.text:
                    safe_run = run.text.replace("¥", "Yen")
                    safe_old = old_text.replace("¥", "Yen")
                    safe_new = new_text.replace("¥", "Yen")
                    logger.info(f"    Updating run: '{safe_run}' -> replacing '{safe_old}' with '{safe_new}'")
                    run.text = run.text.replace(old_text, new_text)
                    
    if shape.has_table:
        table = shape.table
        for row in table.rows:
            for cell in row.cells:
                if old_text in cell.text:
                    safe_cell = cell.text.replace("¥", "Yen")
                    safe_old = old_text.replace("¥", "Yen")
                    safe_new = new_text.replace("¥", "Yen")
                    logger.info(f"    Updating table cell: '{safe_cell}' -> replacing '{safe_old}' with '{safe_new}'")
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)

def main():
    parser = argparse.ArgumentParser(description="Refresh presentation slide deck numbers")
    parser.add_argument("ticker", type=str, help="Stock ticker")
    parser.add_argument("--pptx-path", type=str, help="Path to PPTX file (optional)")
    parser.add_argument("--replacements", type=str, help="Comma-separated replacement mapping (e.g. 'old1:new1,old2:new2')")
    parser.add_argument("--outdir", type=str, default="./out", help="Base output directory")
    args = parser.parse_args()
    
    ticker_str = args.ticker.strip()
    ticker_dir = find_ticker_dir(args.outdir, ticker_str)
    
    if not os.path.exists(ticker_dir):
        logger.error(f"Directory {ticker_dir} does not exist.")
        sys.exit(1)
        
    pptx_path = args.pptx_path
    if not pptx_path:
        # ディレクトリ内から *.pptx を自動探索
        analysis_dir = os.path.join(ticker_dir, "analysis")
        pptx_files = glob.glob(os.path.join(analysis_dir, "*.pptx"))
        if pptx_files:
            pptx_path = pptx_files[0]
        else:
            logger.error(f"No presentation files (*.pptx) found in {analysis_dir}. Please specify --pptx-path.")
            sys.exit(1)
            
    if not os.path.exists(pptx_path):
        logger.error(f"Presentation not found: {pptx_path}")
        sys.exit(1)
        
    logger.info(f"Loading presentation for deck refresh: {pptx_path}...")
    prs = Presentation(pptx_path)
    
    # 置換マッピングの構築
    replacements = {}
    if args.replacements:
        pairs = args.replacements.split(",")
        for pair in pairs:
            if ":" in pair:
                k, v = pair.split(":", 1)
                replacements[k.strip()] = v.strip()
    else:
        # デフォルトはキオクシアのデモ用マッピング（後方互換性用）
        if "285A" in ticker_str:
            replacements = {
                "¥1,850.0 Bn": "¥1,920.0 Bn",
                "¥820.0 Bn": "¥880.0 Bn",
                "44.3%": "45.8%"
            }
        else:
            logger.warning("No replacements specified. Slide content will not be modified.")
            
    if replacements:
        print("=== Deck Refresh Execution Plan ===")
        for k, v in replacements.items():
            safe_k = k.replace("¥", "Yen")
            safe_v = v.replace("¥", "Yen")
            print(f"  Mapping: '{safe_k}' -> '{safe_v}'")
            
        for slide_idx, slide in enumerate(prs.slides, 1):
            logger.info(f"  Scanning Slide {slide_idx}...")
            for shape in slide.shapes:
                for old_val, new_val in replacements.items():
                    refresh_text(shape, old_val, new_val)
                    
        prs.save(pptx_path)
        logger.info(f"Deck refresh complete. Saved updated presentation back to {pptx_path}")
        print("===================================\n")
    else:
        logger.info("Nothing to refresh.")

if __name__ == "__main__":
    main()
