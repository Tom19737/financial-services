import openpyxl
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os
import sys
import json
import argparse
from datetime import datetime
from utils import find_ticker_dir, normalize_ticker, get_latest_financial_data, setup_logging, sanitize_folder_name

logger = setup_logging("generate_pitch")

def main():
    parser = argparse.ArgumentParser(description="Generate investor pitch presentation slide deck")
    parser.add_argument("ticker", type=str, help="Stock ticker")
    parser.add_argument("--outdir", type=str, default="./out", help="Base output directory")
    args = parser.parse_args()
    
    ticker_str = args.ticker.strip()
    ticker_dir = find_ticker_dir(args.outdir, ticker_str)
    
    if not os.path.exists(ticker_dir):
        logger.error(f"Directory {ticker_dir} does not exist.")
        sys.exit(1)
        
    try:
        ticker_data = get_latest_financial_data(ticker_dir, ticker_str)
    except FileNotFoundError as e:
        logger.error(str(e))
        sys.exit(1)
        
    summary_data = {}
    sum_path = os.path.join(ticker_dir, "market_data", "summary.json")
    if os.path.exists(sum_path):
        try:
            with open(sum_path, "r", encoding="utf-8") as f:
                summary_data = json.load(f)
        except Exception as e:
            logger.warning(f"Failed to read summary.json: {e}")
            
    prs = Presentation()
    
    # プレゼンテーションサイズを 16:9 ワイドスクリーンに設定
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # 配色設定 (ダークネイビーとゴールド／イエローをベースにした高級感あるテーマ)
    c_navy = RGBColor(27, 38, 59)
    c_gold = RGBColor(212, 175, 55)
    c_dark = RGBColor(30, 30, 30)
    c_gray = RGBColor(120, 120, 120)
    c_white = RGBColor(255, 255, 255)
    
    font_title = "Outfit"
    font_body = "Outfit"

    # ヘルパー関数: スライドにタイトルと背景色を追加
    def set_slide_base(slide, title_text, dark_bg=False):
        background = slide.background
        fill = background.fill
        if dark_bg:
            fill.solid()
            fill.fore_color.rgb = c_navy
        else:
            fill.solid()
            fill.fore_color.rgb = c_white
            
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = font_title
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = c_white if dark_bg else c_navy

    # Slide 1: Title Slide (Dark Background)
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_base(slide1, "", dark_bg=True)
    
    # 中央のメインタイトル
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.2), Inches(11.33), Inches(2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = f"{ticker_data['name']} ({ticker_data['ticker']})"
    p.font.name = font_title
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = c_gold
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = "投資家向けピッチプレゼンテーション"
    p2.font.name = font_title
    p2.font.size = Pt(24)
    p2.font.bold = True
    p2.font.color.rgb = c_white
    
    # 日付と署名
    meta_box = slide1.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11.33), Inches(1))
    tf_meta = meta_box.text_frame
    p_meta = tf_meta.paragraphs[0]
    p_meta.alignment = PP_ALIGN.CENTER
    p_meta.text = f"{datetime.now().strftime('%Y年%m月%d日')} | アナリスト作成資料"
    p_meta.font.name = font_body
    p_meta.font.size = Pt(14)
    p_meta.font.color.rgb = c_gray

    # Slide 2: Executive Summary (Light Background)
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_base(slide2, "I. エグゼクティブ・サマリー")
    
    # メイン提案ボックス
    prop_box = slide2.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(6), Inches(4.5))
    tf_prop = prop_box.text_frame
    tf_prop.word_wrap = True
    p_prop = tf_prop.paragraphs[0]
    p_prop.text = "投資判断："
    p_prop.font.name = font_title
    p_prop.font.size = Pt(18)
    p_prop.font.bold = True
    p_prop.font.color.rgb = c_navy
    
    p_rec = tf_prop.add_paragraph()
    p_rec.text = "買い (ロング)"
    p_rec.font.name = font_title
    p_rec.font.size = Pt(36)
    p_rec.font.bold = True
    p_rec.font.color.rgb = RGBColor(4, 120, 87) # Emerald Green
    
    is_jpy = ticker_data["currency"] == "JPY"
    currency_symbol = "¥" if is_jpy else "$"
    div_factor = 1e9 if is_jpy else 1e6
    unit_suffix = "十億円" if is_jpy else "百万ドル"
    
    rev_val = ticker_data["revenue"] / div_factor
    eb_val = ticker_data["ebitda"] / div_factor
    ebitda_margin = (ticker_data["ebitda"] / ticker_data["revenue"]) if ticker_data["revenue"] else 0.0
    shares_m = ticker_data["shares_outstanding"] / 1e6
    current_price = ticker_data["current_price"]
    
    p_target = tf_prop.add_paragraph()
    # 目標株価の簡易算定 (DCFモデルがあればそちらから読み出せるが、ない場合は仮に30%プレミアム)
    implied_target = current_price * 1.3
    dcf_file = os.path.join(ticker_dir, "analysis", f"dcf_{ticker_data['ticker']}.xlsx")
    if os.path.exists(dcf_file):
        try:
            wb = openpyxl.load_workbook(dcf_file, data_only=True)
            ws = wb["DCFバリュエーション"] if "DCFバリュエーション" in wb.sheetnames else (wb["DCF Valuation"] if "DCF Valuation" in wb.sheetnames else wb.active)
            val_price = ws["B47"].value
            if val_price is not None:
                implied_target = float(val_price)
            wb.close()
        except Exception as e:
            logger.warning(f"Failed to read target price from DCF model for pitch: {e}")
            
    upside = ((implied_target - current_price) / current_price * 100) if current_price else 0.0
    p_target.text = f"目標株価: {currency_symbol}{implied_target:,.1f}\n現在株価: {currency_symbol}{current_price:,.1f} (上昇余地 {upside:+.1f}%)"
    p_target.font.name = font_body
    p_target.font.size = Pt(16)
    p_target.font.color.rgb = c_dark
    
    p_desc = tf_prop.add_paragraph()
    summary_desc = summary_data.get("description")
    if summary_desc:
        p_desc.text = f"\n{summary_desc}"
    else:
        sector = summary_data.get("sector") or "テクノロジー"
        industry = summary_data.get("industry") or "半導体"
        p_desc.text = f"\n{ticker_data['name']}は、{sector}セクターの{industry}業界におけるリーディング企業です。堅実な財務業績を維持しており、高い成長機会と強固なキャッシュ創出力を備えています。"
        
    p_desc.font.name = font_body
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = c_dark
    
    # 統計指標テーブル
    rows, cols = 5, 2
    left, top, width, height = Inches(7), Inches(1.8), Inches(5.8), Inches(4)
    table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    table.columns[0].width = Inches(3.3)
    table.columns[1].width = Inches(2.5)
    
    table_data = [
        ("FY25E 売上高", f"{currency_symbol}{rev_val:,.1f} {unit_suffix}"),
        ("FY25E EBITDA", f"{currency_symbol}{eb_val:,.1f} {unit_suffix}"),
        ("FY25E EBITDAマージン", f"{ebitda_margin:.1%}"),
        ("発行済株式数", f"{shares_m:,.1f} 百万株"),
        ("現在株価", f"{currency_symbol}{current_price:,.1f}")
    ]
    for row_idx, (k, v) in enumerate(table_data):
        cell_k = table.cell(row_idx, 0)
        cell_v = table.cell(row_idx, 1)
        cell_k.text = k
        cell_v.text = v
        for cell in [cell_k, cell_v]:
            for p in cell.text_frame.paragraphs:
                p.font.name = font_body
                p.font.size = Pt(14)
                p.font.bold = True if row_idx == 0 else False
                p.font.color.rgb = c_dark
                p.alignment = PP_ALIGN.LEFT if cell == cell_k else PP_ALIGN.RIGHT

    # Slide 3: Key Investment Pillars (Light Background)
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_base(slide3, "II. 主要な投資ポイント")
    
    # 3つのピラーを列ボックスで配置
    widths = Inches(3.8)
    gap = Inches(0.4)
    
    summary_pillars = summary_data.get("pillars")
    if summary_pillars:
        pillars = [(item[0], item[1]) for item in summary_pillars]
    else:
        industry = summary_data.get("industry") or "業界"
        pillars = [
            ("1. 業界におけるリーダーシップ", f"{ticker_data['name']}は、{industry}分野において圧倒的な市場ポジションを確立しています。強いブランド認知度と独自の技術力は、強固な競争優位性（経済的な堀）となっています。"),
            ("2. 卓越したオペレーショナル・エフィシエンシー", "優れたサプライチェーン管理と製造効率により、高い営業レバレッジを実現しています。継続的な最適化により、競合他社を上回る高いEBITDAマージンを維持しています。"),
            ("3. 財務の柔軟性", "十分な手元資金と管理可能な負債水準を備えた、健全なバランスシートを有しています。これにより、M&Aや研究開発、株主還元に対する戦略的選択肢がもたらされます。")
        ]
        
    for idx, (title, text) in enumerate(pillars):
        x_pos = Inches(0.5) + idx * (widths + gap)
        col_box = slide3.shapes.add_textbox(x_pos, Inches(1.8), widths, Inches(4.5))
        tf_col = col_box.text_frame
        tf_col.word_wrap = True
        
        p_title = tf_col.paragraphs[0]
        p_title.text = title
        p_title.font.name = font_title
        p_title.font.size = Pt(18)
        p_title.font.bold = True
        p_title.font.color.rgb = c_navy
        
        p_body = tf_col.add_paragraph()
        p_body.text = "\n" + text
        p_body.font.name = font_body
        p_body.font.size = Pt(13)
        p_body.font.color.rgb = c_dark
        
    target_path = os.path.join(ticker_dir, "analysis")
    os.makedirs(target_path, exist_ok=True)
    
    clean_name = sanitize_folder_name(ticker_data["name"])
    today_str = datetime.now().strftime("%Y%m%d")
    out_file = os.path.join(target_path, f"{clean_name}_Pitch_{today_str}.pptx")
    
    prs.save(out_file)
    logger.info(f"Presentation saved to {out_file}")

if __name__ == "__main__":
    main()
